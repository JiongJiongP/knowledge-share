"""Convert Reveal.js HTML presentation to PPTX."""
from html.parser import HTMLParser
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import os

HTML_FILE = "index.html"
OUTPUT_FILE = "presentation.pptx"

# ── Color scheme (dark theme) ──
BG_DARK = RGBColor(0x19, 0x19, 0x19)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x64, 0xB5, 0xF6)
GREEN = RGBColor(0x43, 0xE9, 0x7B)
ORANGE = RGBColor(0xFF, 0xB7, 0x4D)
PINK = RGBColor(0xF0, 0x62, 0x92)
RED = RGBColor(0xEF, 0x53, 0x50)
GRAY = RGBColor(0x78, 0x90, 0x9C)
LIGHT_BLUE = RGBColor(0x90, 0xCA, 0xF9)


class SlideExtractor(HTMLParser):
    """Extract slide content from Reveal.js HTML."""

    def __init__(self):
        super().__init__()
        self.slides = []
        self._current_slide = None
        self._in_slides_div = False
        self._in_section = False
        self._section_depth = 0
        self._skip_tags = {"style", "script", "aside"}
        self._skip_depth = 0
        self._tag_stack = []
        self._text_buffer = []

    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)

        if self._skip_depth > 0:
            self._skip_depth += 1
            return

        if tag in self._skip_tags:
            self._skip_depth = 1
            return

        # Track entry into slides container
        if tag == "div" and "slides" in attrs_dict.get("class", ""):
            self._in_slides_div = True
            return

        if tag == "section" and self._in_slides_div:
            self._in_section = True
            self._section_depth += 1
            if self._section_depth == 1:
                self._current_slide = {"type": "content", "elements": []}
                self._text_buffer = []
            return

        if self._in_section:
            self._tag_stack.append(tag)
            if tag in ("h1", "h2", "h3", "h4"):
                self._text_buffer.append(f"<{tag}>")
            elif tag == "li":
                self._text_buffer.append("<li>")
            elif tag == "p":
                self._text_buffer.append("<p>")
            elif tag == "strong" or tag == "b":
                self._text_buffer.append("<b>")
            elif tag == "em":
                self._text_buffer.append("<i>")
            elif tag == "br":
                self._text_buffer.append("\n")
            elif tag == "code":
                self._text_buffer.append("<code>")
            elif tag == "span":
                cls = attrs_dict.get("class", "")
                if "accent-blue" in cls:
                    self._text_buffer.append("<blue>")
                elif "accent-green" in cls:
                    self._text_buffer.append("<green>")
                elif "accent-orange" in cls:
                    self._text_buffer.append("<orange>")
                elif "accent-pink" in cls:
                    self._text_buffer.append("<pink>")
                elif "accent-red" in cls:
                    self._text_buffer.append("<red>")

    def handle_endtag(self, tag):
        if self._skip_depth > 0:
            self._skip_depth -= 1
            return

        if tag == "div" and self._in_slides_div:
            return

        if tag == "section" and self._in_section:
            self._section_depth -= 1
            if self._section_depth == 0:
                # End of top-level section
                raw = "".join(self._text_buffer).strip()
                self._current_slide["raw_text"] = raw
                self.slides.append(self._current_slide)
                self._current_slide = None
                self._text_buffer = []
            return

        if self._in_section:
            if tag in ("h1", "h2", "h3", "h4"):
                self._text_buffer.append(f"</{tag}>\n")
            elif tag == "li":
                self._text_buffer.append("</li>\n")
            elif tag == "p":
                self._text_buffer.append("</p>\n")
            elif tag == "strong" or tag == "b":
                self._text_buffer.append("</b>")
            elif tag == "em":
                self._text_buffer.append("</i>")
            elif tag == "code":
                self._text_buffer.append("</code>")
            elif tag == "span":
                pass  # handled via custom markers
            if self._tag_stack and self._tag_stack[-1] == tag:
                self._tag_stack.pop()

    def handle_data(self, data):
        if self._skip_depth > 0:
            return
        if self._in_section:
            self._text_buffer.append(data)


def parse_styled_text(raw_text):
    """Parse simplified markup into list of (text, style) tuples.
    Styles: h1, h2, h3, h4, li, p, code, b, i, blue, green, orange, pink, red
    """
    import re

    parts = []
    # Tokenize
    tokens = re.split(
        r'(</?h[1-4]>|</?li>|</?p>|</?b>|</?i>|</?code>|'
        r'</?blue>|</?green>|</?orange>|</?pink>|</?red>)',
        raw_text,
    )

    current_style = "p"
    for token in tokens:
        token = token.strip()
        if not token:
            continue
        if token == "<h1>":
            current_style = "h1"
        elif token == "<h2>":
            current_style = "h2"
        elif token == "<h3>":
            current_style = "h3"
        elif token == "<h4>":
            current_style = "h4"
        elif token == "<p>":
            current_style = "p"
        elif token == "<li>":
            current_style = "li"
        elif token == "<b>":
            current_style = "b"
        elif token == "<i>":
            current_style = "i"
        elif token == "<code>":
            current_style = "code"
        elif token == "<blue>":
            current_style = "blue"
        elif token == "<green>":
            current_style = "green"
        elif token == "<orange>":
            current_style = "orange"
        elif token == "<pink>":
            current_style = "pink"
        elif token == "<red>":
            current_style = "red"
        elif token.startswith("</"):
            current_style = "p"
        else:
            parts.append((token, current_style))
    return parts


def add_dark_background(slide):
    """Set dark background for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def create_slide_content(slide, raw_text):
    """Parse slide text and add it to the slide."""
    if not raw_text:
        return

    parts = parse_styled_text(raw_text)

    # Group parts by their style for contiguous runs
    left = Inches(0.8)
    top = Inches(0.5)
    width = Inches(8.4)
    height = Inches(6.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Remove default empty paragraph
    if tf.paragraphs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    first = True
    for text, style in parts:
        text = text.strip()
        if not text:
            continue

        if style in ("h1", "h2", "h3", "h4", "li"):
            # New paragraph for block elements
            if not first or (first and p.text):
                p = tf.add_paragraph()
            first = False

        run = p.add_run()
        run.text = text

        # Apply styling
        if style == "h1":
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        elif style == "h2":
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = BLUE
        elif style == "h3":
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = LIGHT_BLUE
        elif style == "h4":
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
        elif style == "li":
            run.font.size = Pt(14)
            run.font.color.rgb = WHITE
            p.level = 0
            # Add bullet character
            run.text = "• " + text
        elif style == "code":
            run.font.size = Pt(12)
            run.font.name = "Consolas"
            run.font.color.rgb = ORANGE
        elif style == "b":
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = WHITE
        elif style == "blue":
            run.font.color.rgb = BLUE
            run.font.size = Pt(14)
        elif style == "green":
            run.font.color.rgb = GREEN
            run.font.size = Pt(14)
        elif style == "orange":
            run.font.color.rgb = ORANGE
            run.font.size = Pt(14)
        elif style == "pink":
            run.font.color.rgb = PINK
            run.font.size = Pt(14)
        elif style == "red":
            run.font.color.rgb = RED
            run.font.size = Pt(14)
        else:
            run.font.size = Pt(14)
            run.font.color.rgb = WHITE


def create_title_slide(prs, raw_text):
    """Create a title/cover slide."""
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)

    if not raw_text:
        return

    parts = parse_styled_text(raw_text)
    full_text = " ".join(t for t, s in parts if s in ("h1",) or (s == "p" and t))

    # Main title
    if any(s == "h1" for _, s in parts):
        title_text = " ".join(t for t, s in parts if s == "h1")
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title_text
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = BLUE

    # Subtitle
    subtitle_parts = [t for t, s in parts if s == "p" and not t.startswith("•")]
    if subtitle_parts:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, sub in enumerate(subtitle_parts):
            if i > 0:
                tf2.add_paragraph()
            p2 = tf2.paragraphs[-1]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = sub
            run2.font.size = Pt(16)
            run2.font.color.rgb = GRAY


def main():
    # Read HTML
    with open(HTML_FILE, "r", encoding="utf-8") as f:
        html = f.read()

    # Extract slides
    extractor = SlideExtractor()
    extractor.feed(html)

    print(f"Found {len(extractor.slides)} slides")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(extractor.slides):
        raw = slide_data.get("raw_text", "")

        # Detect title slide
        is_title = i == 0

        if is_title:
            create_title_slide(prs, raw)
        else:
            slide_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(slide_layout)
            add_dark_background(slide)
            create_slide_content(slide, raw)

        # Add slide number
        if not is_title:
            txBox = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.4))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            run = p.add_run()
            run.text = str(i + 1)
            run.font.size = Pt(8)
            run.font.color.rgb = GRAY

    # Save
    prs.save(OUTPUT_FILE)
    print(f"Saved to {OUTPUT_FILE}")
    print(f"File size: {os.path.getsize(OUTPUT_FILE) / 1024:.1f} KB")


if __name__ == "__main__":
    main()
